from pptx import Presentation

def scan_pptx(file_path):

    findings = []

    try:

        prs = Presentation(file_path)

        findings.append({
            "type": "Slides",
            "value": len(prs.slides)
        })

        for slide in prs.slides:

            for shape in slide.shapes:

                if hasattr(shape, "text"):

                    findings.append({
                        "type": "Text Found",
                        "value": shape.text[:100]
                    })

    except Exception as e:

        findings.append({
            "type": "Error",
            "value": str(e)
        })

    return findings